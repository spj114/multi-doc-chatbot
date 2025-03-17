import sys
import os
import io
from PyQt5.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout, 
                           QHBoxLayout, QPushButton, QLabel, QTextEdit, 
                           QListWidget, QFileDialog, QLineEdit, QProgressBar,
                           QMessageBox, QListWidgetItem)
from PyQt5.QtCore import Qt, QThread, pyqtSignal
from PyQt5.QtGui import QIcon, QFont

# Import libraries from your original code
from PyPDF2 import PdfReader
import fitz  # PyMuPDF for image extraction
from PIL import Image
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
import docx  # for Word documents
import pandas as pd  # for Excel files
from pptx import Presentation  # for PowerPoint files

# Load environment variables
load_dotenv()
api_key = os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=api_key)

# Processing thread class
class ProcessingThread(QThread):
    progress_update = pyqtSignal(int, str)
    finished = pyqtSignal(bool)
    
    def __init__(self, files):
        super().__init__()
        self.files = files
        self.text = ""
        self.images = []
    
    def run(self):
        try:
            total_files = len(self.files)
            for i, file_path in enumerate(self.files):
                self.progress_update.emit(int((i / total_files) * 50), f"Processing {os.path.basename(file_path)}")
                
                file_extension = os.path.splitext(file_path)[1].lower()
                
                try:
                    # Process based on file type
                    if file_extension == '.pdf':
                        # PDF processing
                        self.progress_update.emit(int((i / total_files) * 50), f"Extracting text from PDF: {os.path.basename(file_path)}")
                        pdf_reader = PdfReader(file_path)
                        for page in pdf_reader.pages:
                            self.text += page.extract_text() or ""
                        
                        # Extract images from PDF
                        self.progress_update.emit(int((i / total_files) * 50), f"Extracting images from PDF: {os.path.basename(file_path)}")
                        doc_pdf = fitz.open(file_path)
                        for page in doc_pdf:
                            for img_index, img in enumerate(page.get_images(full=True)):
                                xref = img[0]
                                base_image = doc_pdf.extract_image(xref)
                                self.images.append(base_image["image"])
                        doc_pdf.close()
                        
                    elif file_extension == '.docx':
                        self.progress_update.emit(int((i / total_files) * 50), f"Processing Word document: {os.path.basename(file_path)}")
                        self.text += self.extract_text_from_word(file_path)
                        
                    elif file_extension in ['.xlsx', '.xls']:
                        self.progress_update.emit(int((i / total_files) * 50), f"Processing Excel file: {os.path.basename(file_path)}")
                        self.text += self.extract_text_from_excel(file_path)
                        
                    elif file_extension in ['.pptx', '.ppt']:
                        self.progress_update.emit(int((i / total_files) * 50), f"Processing PowerPoint file: {os.path.basename(file_path)}")
                        self.text += self.extract_text_from_ppt(file_path)
                        
                    elif file_extension in ['.jpg', '.jpeg', '.png']:
                        # For image files, directly append to images list
                        self.progress_update.emit(int((i / total_files) * 50), f"Processing image file: {os.path.basename(file_path)}")
                        with open(file_path, 'rb') as img_file:
                            self.images.append(img_file.read())
                        
                        # Also analyze the image immediately with Gemini for text extraction
                        img = Image.open(file_path)
                        model = genai.GenerativeModel("gemini-1.5-flash")
                        response = model.generate_content(["Describe this image and extract any visible text:", img])
                        self.text += response.text + "\n" if response and response.text else ""
                        
                except Exception as e:
                    self.progress_update.emit(int((i / total_files) * 50), f"Error processing {os.path.basename(file_path)}: {str(e)}")
            
            # Analyze images with Gemini
            self.progress_update.emit(50, "Analyzing images with Gemini...")
            image_text = self.analyze_images_with_gemini(self.images)
            combined_text = self.text + "\n" + image_text
            
            # Process text chunks
            self.progress_update.emit(70, "Processing text chunks...")
            text_chunks = self.get_text_chunks(combined_text)
            
            # Create vector store
            self.progress_update.emit(85, "Creating vector store...")
            self.get_vector_store(text_chunks)
            
            self.progress_update.emit(100, "Processing complete!")
            self.finished.emit(True)
            
        except Exception as e:
            self.progress_update.emit(0, f"Processing failed: {str(e)}")
            self.finished.emit(False)
    
    def extract_text_from_word(self, doc_file):
        doc = docx.Document(doc_file)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text

    def extract_text_from_excel(self, excel_file):
        df = pd.read_excel(excel_file)
        return df.to_string()

    def extract_text_from_ppt(self, ppt_file):
        prs = Presentation(ppt_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    
    def analyze_images_with_gemini(self, images):
        extracted_text = ""
        model = genai.GenerativeModel("gemini-1.5-flash")
        prompt = "Analyze this image and describe its contents in detail, including any visible text, diagrams, or visual elements."

        for i, img_data in enumerate(images):
            try:
                self.progress_update.emit(50 + int((i / len(images)) * 10), f"Analyzing image {i+1}/{len(images)}...")
                img = Image.open(io.BytesIO(img_data))
                response = model.generate_content([prompt, img])
                extracted_text += response.text + "\n" if response and response.text else ""
            except Exception as e:
                self.progress_update.emit(50 + int((i / len(images)) * 10), f"Error processing image {i+1}: {str(e)}")

        return extracted_text
    
    def get_text_chunks(self, text):
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
        chunks = text_splitter.split_text(text)
        return chunks

    def get_vector_store(self, text_chunks):
        embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
        vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
        vector_store.save_local("faiss_index")

# Chat query thread
class ChatQueryThread(QThread):
    response_ready = pyqtSignal(str)
    
    def __init__(self, question):
        super().__init__()
        self.question = question
    
    def run(self):
        try:
            embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
            new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
            docs = new_db.similarity_search(self.question)
            
            prompt_template = """
            Answer the question as detailed as possible from the provided context, including text and extracted image data.
            If the answer is not in the context, say "answer is not available in the context". Do not guess.

            Context:
            {context}
            
            Question:
            {question}

            Answer:
            """
            
            model = ChatGoogleGenerativeAI(model="gemini-1.5-pro", temperature=0.3)
            prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
            chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
            
            response = chain({"input_documents": docs, "question": self.question}, return_only_outputs=True)
            self.response_ready.emit(response["output_text"])
            
        except Exception as e:
            self.response_ready.emit(f"Error: {str(e)}")

# Main application class
class DocumentChatApp(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("Document Chat")
        self.setGeometry(100, 100, 1200, 700)
        self.setStyleSheet("""
            QMainWindow {
                background-color: #f0f0f0;
            }
            QPushButton {
                background-color: #4285f4;
                color: white;
                border: none;
                padding: 8px 16px;
                border-radius: 4px;
                font-weight: bold;
            }
            QPushButton:hover {
                background-color: #3367d6;
            }
            QPushButton:disabled {
                background-color: #cccccc;
            }
            QTextEdit, QListWidget, QLineEdit {
                border: 1px solid #dcdcdc;
                border-radius: 4px;
                padding: 4px;
            }
            QLabel {
                font-weight: bold;
            }
        """)
        
        # Main layout
        main_widget = QWidget()
        main_layout = QHBoxLayout()
        main_widget.setLayout(main_layout)
        self.setCentralWidget(main_widget)
        
        # Left panel (document management)
        left_panel = QWidget()
        left_layout = QVBoxLayout()
        left_panel.setLayout(left_layout)
        left_panel.setFixedWidth(350)
        
        app_title = QLabel("Document Chat")
        app_title.setFont(QFont("Arial", 16, QFont.Bold))
        app_title.setAlignment(Qt.AlignCenter)
        app_title.setStyleSheet("color: #4285f4; margin: 10px;")
        
        upload_btn = QPushButton("Upload Documents")
        upload_btn.setIcon(QIcon.fromTheme("document-open"))
        upload_btn.clicked.connect(self.upload_documents)
        
        self.file_list = QListWidget()
        self.file_list.setStyleSheet("QListWidget { background-color: white; }")
        
        self.process_btn = QPushButton("Process Documents")
        self.process_btn.setIcon(QIcon.fromTheme("system-run"))
        self.process_btn.clicked.connect(self.process_documents)
        self.process_btn.setEnabled(False)
        
        self.progress_bar = QProgressBar()
        self.progress_bar.setRange(0, 100)
        self.progress_bar.setValue(0)
        self.progress_bar.setTextVisible(True)
        self.progress_bar.setFormat("%p% %v")
        
        self.status_label = QLabel("Ready")
        self.status_label.setAlignment(Qt.AlignCenter)
        self.status_label.setStyleSheet("color: #666;")
        
        left_layout.addWidget(app_title)
        left_layout.addWidget(upload_btn)
        left_layout.addWidget(QLabel("Uploaded Documents:"))
        left_layout.addWidget(self.file_list)
        left_layout.addWidget(self.process_btn)
        left_layout.addWidget(self.progress_bar)
        left_layout.addWidget(self.status_label)
        left_layout.addStretch()
        
        # Right panel (chat interface)
        right_panel = QWidget()
        right_layout = QVBoxLayout()
        right_panel.setLayout(right_layout)
        
        chat_title = QLabel("Chat with your Documents")
        chat_title.setFont(QFont("Arial", 14, QFont.Bold))
        chat_title.setAlignment(Qt.AlignCenter)
        
        self.chat_history = QTextEdit()
        self.chat_history.setReadOnly(True)
        self.chat_history.setStyleSheet("""
            QTextEdit {
                background-color: white;
                font-size: 12pt;
            }
        """)
        
        input_layout = QHBoxLayout()
        self.chat_input = QLineEdit()
        self.chat_input.setPlaceholderText("Ask a question about your documents...")
        self.chat_input.returnPressed.connect(self.send_message)
        self.chat_input.setEnabled(False)
        
        send_btn = QPushButton("Send")
        send_btn.setIcon(QIcon.fromTheme("mail-send"))
        send_btn.clicked.connect(self.send_message)
        send_btn.setEnabled(False)
        self.send_btn = send_btn
        
        input_layout.addWidget(self.chat_input)
        input_layout.addWidget(send_btn)
        
        right_layout.addWidget(chat_title)
        right_layout.addWidget(self.chat_history)
        right_layout.addLayout(input_layout)
        
        # Add panels to main layout
        main_layout.addWidget(left_panel)
        main_layout.addWidget(right_panel, 1)  # 1 is the stretch factor
        
        # Initialize variables
        self.files = []
        self.is_processed = False
        
        # Show welcome message
        self.chat_history.append("<div style='color: #4285f4; font-weight: bold;'>Welcome to Document Chat!</div>")
        self.chat_history.append("<div style='color: #666;'>Upload and process documents to start chatting with them.</div>")
    
    def upload_documents(self):
        filenames, _ = QFileDialog.getOpenFileNames(
            self, "Upload Documents", "", 
            "Documents (*.pdf *.docx *.xlsx *.xls *.pptx *.ppt *.jpg *.jpeg *.png)"
        )
        
        if filenames:
            self.files = filenames  # Replace existing files
            self.file_list.clear()
            for file in self.files:
                item = QListWidgetItem(os.path.basename(file))
                # Set icon based on file type
                file_extension = os.path.splitext(file)[1].lower()
                if file_extension == '.pdf':
                    item.setIcon(QIcon.fromTheme("application-pdf"))
                elif file_extension in ['.docx', '.doc']:
                    item.setIcon(QIcon.fromTheme("application-msword"))
                elif file_extension in ['.xlsx', '.xls']:
                    item.setIcon(QIcon.fromTheme("application-vnd.ms-excel"))
                elif file_extension in ['.pptx', '.ppt']:
                    item.setIcon(QIcon.fromTheme("application-vnd.ms-powerpoint"))
                elif file_extension in ['.jpg', '.jpeg', '.png']:
                    item.setIcon(QIcon.fromTheme("image"))
                else:
                    item.setIcon(QIcon.fromTheme("text-plain"))
                
                self.file_list.addItem(item)
            
            self.process_btn.setEnabled(True)
            self.status_label.setText(f"{len(self.files)} documents ready for processing")
    
    def process_documents(self):
        if not self.files:
            QMessageBox.warning(self, "No Documents", "Please upload documents first.")
            return
        
        # Reset UI
        self.progress_bar.setValue(0)
        self.process_btn.setEnabled(False)
        self.is_processed = False
        self.chat_input.setEnabled(False)
        self.send_btn.setEnabled(False)
        
        # Start processing thread
        self.processing_thread = ProcessingThread(self.files)
        self.processing_thread.progress_update.connect(self.update_progress)
        self.processing_thread.finished.connect(self.on_processing_finished)
        self.processing_thread.start()
        
        # Show processing status
        self.chat_history.append("<div style='color: #FF9800; font-weight: bold;'>System: Processing documents...</div>")
    
    def update_progress(self, value, message):
        self.progress_bar.setValue(value)
        self.status_label.setText(message)
    
    def on_processing_finished(self, success):
        if success:
            self.is_processed = True
            self.chat_input.setEnabled(True)
            self.send_btn.setEnabled(True)
            self.chat_history.append("<div style='color: #4CAF50; font-weight: bold;'>System: Documents processed successfully!</div>")
            self.chat_history.append("<div style='color: #666;'>You can now ask questions about your documents.</div>")
        else:
            self.process_btn.setEnabled(True)
            self.chat_history.append("<div style='color: #F44336; font-weight: bold;'>System: Error processing documents.</div>")
            self.chat_history.append("<div style='color: #666;'>Please check the console for details and try again.</div>")
    
    def send_message(self):
        if not self.is_processed:
            self.chat_history.append("<div style='color: #FF9800; font-weight: bold;'>System: Please process documents first.</div>")
            return
        
        question = self.chat_input.text().strip()
        if not question:
            return
        
        # Display user message
        self.chat_history.append(f"<div style='color: #4285f4; font-weight: bold;'>You:</div> {question}")
        self.chat_input.clear()
        self.chat_input.setEnabled(False)
        self.send_btn.setEnabled(False)
        
        # Show "thinking" indicator
        self.status_label.setText("Thinking...")
        
        # Start query thread
        self.query_thread = ChatQueryThread(question)
        self.query_thread.response_ready.connect(self.display_response)
        self.query_thread.start()
    
    def display_response(self, response):
        # Display AI response
        self.chat_history.append(f"<div style='color: #0F9D58; font-weight: bold;'>AI:</div> {response}")
        
        # Scroll to bottom
        self.chat_history.verticalScrollBar().setValue(self.chat_history.verticalScrollBar().maximum())
        
        # Re-enable input
        self.chat_input.setEnabled(True)
        self.send_btn.setEnabled(True)
        self.status_label.setText("Ready")

def main():
    # Initialize PyQt application
    app = QApplication(sys.argv)
    
    # Create and show the main window
    window = DocumentChatApp()
    window.show()
    
    # Run the application
    sys.exit(app.exec_())

if __name__ == "__main__":
    main()