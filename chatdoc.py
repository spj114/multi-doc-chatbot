import streamlit as st
from PyPDF2 import PdfReader
import fitz  # PyMuPDF for image extraction
import os
import io
from PIL import Image
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
import docx  # for Word documents
import pandas as pd  # for Excel files
from pptx import Presentation  # for PowerPoint files

load_dotenv()
api_key = os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=api_key)

def extract_text_from_word(doc_file):
    doc = docx.Document(doc_file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def extract_text_from_excel(excel_file):
    df = pd.read_excel(excel_file)
    return df.to_string()

def extract_text_from_ppt(ppt_file):
    prs = Presentation(ppt_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def get_document_text_and_images(docs):
    text = ""
    images = []
    
    for doc in docs:
        # Save the uploaded file temporarily
        with open(doc.name, "wb") as f:
            f.write(doc.read())
        
        file_extension = os.path.splitext(doc.name)[1].lower()
        
        try:
            # Process based on file type
            if file_extension == '.pdf':
                # PDF processing
                pdf_reader = PdfReader(doc.name)
                for page in pdf_reader.pages:
                    text += page.extract_text() or ""
                
                # Extract images from PDF
                doc_pdf = fitz.open(doc.name)
                for page in doc_pdf:
                    for img_index, img in enumerate(page.get_images(full=True)):
                        xref = img[0]
                        base_image = doc_pdf.extract_image(xref)
                        images.append(base_image["image"])
                doc_pdf.close()
                
            elif file_extension == '.docx':
                text += extract_text_from_word(doc.name)
                
            elif file_extension in ['.xlsx', '.xls']:
                text += extract_text_from_excel(doc.name)
                
            elif file_extension in ['.pptx', '.ppt']:
                text += extract_text_from_ppt(doc.name)
                
            elif file_extension in ['.jpg', '.jpeg', '.png']:
                # For image files, directly append to images list
                with open(doc.name, 'rb') as img_file:
                    images.append(img_file.read())
                
                # Also analyze the image immediately with Gemini for text extraction
                img = Image.open(doc.name)
                model = genai.GenerativeModel("gemini-1.5-flash")
                response = model.generate_content(["Describe this image and extract any visible text:", img])
                text += response.text + "\n" if response and response.text else ""
                
        except Exception as e:
            st.error(f"Error processing {doc.name}: {e}")
        
        # Clean up temporary file
        try:
            os.remove(doc.name)
        except:
            pass
    
    return text, images

def analyze_images_with_gemini(images):
    extracted_text = ""
    model = genai.GenerativeModel("gemini-1.5-flash")
    prompt = "Analyze this image and describe its contents in detail, including any visible text, diagrams, or visual elements."

    for img_data in images:
        try:
            img = Image.open(io.BytesIO(img_data))
            response = model.generate_content([prompt, img])
            extracted_text += response.text + "\n" if response and response.text else ""
        except Exception as e:
            st.error(f"Error processing image: {e}")

    return extracted_text

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context, including text and extracted image data.
    If the answer is not in the context, say "answer is not available in the context". Do not guess.

    Context:
    {context}
    
    Question:
    {question}

    Answer:
    """
    
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

def user_input(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)
    
    chain = get_conversational_chain()
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    
    print(response)
    st.write("Reply:", response["output_text"])

def main():
    st.set_page_config("Chat with Documents")
    st.header("Chat with Multiple Document Formats using Gemini Vision💁")

    user_question = st.text_input("Ask a Question from your Documents")
    
    if user_question:
        user_input(user_question)
    
    with st.sidebar:
        st.title("Menu:")
        accepted_types = ["pdf", "docx", "xlsx", "xls", "pptx", "ppt", "jpg", "jpeg", "png"]
        docs = st.file_uploader(
            "Upload your Documents and Click on Submit & Process",
            accept_multiple_files=True,
            type=accepted_types
        )
        
        if not docs:
            st.warning("Please upload documents before processing.")
        elif st.button("Submit & Process"):
            with st.spinner("Processing..."):
                raw_text, images = get_document_text_and_images(docs)
                image_text = analyze_images_with_gemini(images)
                combined_text = raw_text + "\n" + image_text
                text_chunks = get_text_chunks(combined_text)
                get_vector_store(text_chunks)
                st.success("Processing Complete! You can now chat with your documents.")

if __name__ == "__main__":
    main()